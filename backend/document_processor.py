"""Document processing module for extracting text from various file formats.
Images are stored for vision model analysis (no OCR)."""

import os
import sys
import uuid
import json
import base64
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple
import mimetypes
import logging

# Document processing libraries
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import io

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def get_base_data_dir():
    """Get the base data directory path - uses AppData on Windows."""
    env_data_dir = os.environ.get("DATA_DIR")
    if env_data_dir:
        return env_data_dir
    
    if sys.platform == "win32":
        appdata = os.environ.get("APPDATA", os.path.expanduser("~"))
        return os.path.join(appdata, "LLM Council", "data")
    
    return "data"


# Base data directory
BASE_DATA_DIR = get_base_data_dir()

# Storage directory for documents
DOCUMENTS_DIR = os.path.join(BASE_DATA_DIR, "documents")
DOCUMENT_REGISTRY_FILE = os.path.join(BASE_DATA_DIR, "document_registry.json")

# Supported file types
SUPPORTED_EXTENSIONS = {
    '.pdf': 'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.doc': 'application/msword',
    '.txt': 'text/plain',
    '.rtf': 'application/rtf',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.ppt': 'application/vnd.ms-powerpoint',
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.gif': 'image/gif',
    '.webp': 'image/webp',
    '.bmp': 'image/bmp',
    '.tiff': 'image/tiff',
    '.tif': 'image/tiff',
    '.md': 'text/markdown',
}

# Image extensions - will be sent to vision models
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp', '.tiff', '.tif'}

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
MAX_CHUNK_SIZE = 4000  # Characters per chunk for context window


# ===== Directory Management =====

def ensure_directories():
    """Ensure document storage directories exist."""
    Path(DOCUMENTS_DIR).mkdir(parents=True, exist_ok=True)
    Path(BASE_DATA_DIR).mkdir(parents=True, exist_ok=True)


def load_document_registry() -> Dict[str, Any]:
    """Load the document registry from file."""
    ensure_directories()
    if os.path.exists(DOCUMENT_REGISTRY_FILE):
        with open(DOCUMENT_REGISTRY_FILE, 'r') as f:
            return json.load(f)
    return {"documents": {}}


def save_document_registry(registry: Dict[str, Any]):
    """Save the document registry to file."""
    ensure_directories()
    with open(DOCUMENT_REGISTRY_FILE, 'w') as f:
        json.dump(registry, f, indent=2)


# ===== Text Extraction Functions =====

def extract_text_from_pdf(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """
    Extract text from a PDF file.
    
    Returns:
        Tuple of (extracted_text, metadata_dict)
    """
    metadata = {
        "pages": 0,
        "type": "pdf"
    }
    
    try:
        reader = PdfReader(file_path)
        metadata["pages"] = len(reader.pages)
        
        text_parts = []
        
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{page_text}")
        
        return "\n\n".join(text_parts), metadata
    except Exception as e:
        return f"[Error extracting PDF text: {str(e)}]", metadata


def extract_text_from_docx(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Extract text from a DOCX file."""
    metadata = {"paragraphs": 0, "tables": 0, "type": "docx"}
    
    try:
        doc = DocxDocument(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
                metadata["paragraphs"] += 1
        
        # Also extract from tables
        for table in doc.tables:
            metadata["tables"] += 1
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return "\n\n".join(text_parts), metadata
    except Exception as e:
        return f"[Error extracting DOCX text: {str(e)}]", metadata


def extract_text_from_pptx(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Extract text from a PowerPoint file."""
    metadata = {"slides": 0, "type": "pptx"}
    
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            metadata["slides"] += 1
            slide_texts = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if len(slide_texts) > 1:  # More than just the header
                text_parts.append("\n".join(slide_texts))
        
        return "\n\n".join(text_parts), metadata
    except Exception as e:
        return f"[Error extracting PPTX text: {str(e)}]", metadata


def extract_text_from_txt(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Extract text from a plain text file."""
    metadata = {"encoding": "utf-8", "type": "text"}
    
    try:
        # Try UTF-8 first, then fall back to other encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                metadata["encoding"] = encoding
                return content, metadata
            except UnicodeDecodeError:
                continue
        
        # Last resort: read with error replacement
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read(), metadata
    except Exception as e:
        return f"[Error reading text file: {str(e)}]", metadata


def process_image_for_vision(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """
    Process an image file for vision model analysis.
    Returns base64-encoded image data and metadata.
    
    Returns:
        Tuple of (description_text, metadata_dict)
    """
    metadata = {
        "width": 0,
        "height": 0,
        "format": "Unknown",
        "mode": "Unknown",
        "type": "image",
        "is_vision_image": True,
        "base64_data": None
    }
    
    try:
        img = Image.open(file_path)
        metadata["width"] = img.size[0]
        metadata["height"] = img.size[1]
        metadata["format"] = img.format or "Unknown"
        metadata["mode"] = img.mode
        
        # Convert to RGB if necessary for consistent processing
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')
        
        # Resize if too large (max 2048px on longest side for API limits)
        max_size = 2048
        if max(img.size) > max_size:
            ratio = max_size / max(img.size)
            new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
            img = img.resize(new_size, Image.LANCZOS)
            metadata["resized_to"] = new_size
        
        # Encode to base64
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        base64_data = base64.b64encode(buffer.getvalue()).decode('utf-8')
        metadata["base64_data"] = base64_data
        
        # Return a placeholder text - the actual image will be sent to vision model
        description = f"[Image for Vision Analysis: {metadata['format']} format, {metadata['width']}x{metadata['height']} pixels]"
        return description, metadata
        
    except Exception as e:
        return f"[Error processing image: {str(e)}]", metadata


def extract_text_from_file(file_path: str, extension: str) -> Tuple[str, Dict[str, Any]]:
    """
    Extract text from a file based on its extension.
    Images are processed for vision model analysis.
    
    Args:
        file_path: Path to the file
        extension: File extension (e.g., '.pdf')
    
    Returns:
        Tuple of (extracted_text, metadata_dict)
    """
    ext = extension.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif ext in ['.txt', '.rtf', '.md']:
        return extract_text_from_txt(file_path)
    elif ext in IMAGE_EXTENSIONS:
        return process_image_for_vision(file_path)
    else:
        return f"[Unsupported file format: {ext}]", {}


# ===== Text Chunking =====

def chunk_text(text: str, max_chunk_size: int = MAX_CHUNK_SIZE) -> List[str]:
    """Split text into chunks for LLM context windows."""
    if len(text) <= max_chunk_size:
        return [text]
    
    chunks = []
    current_chunk = ""
    
    # Split by paragraphs first
    paragraphs = text.split("\n\n")
    
    for para in paragraphs:
        if len(current_chunk) + len(para) + 2 <= max_chunk_size:
            if current_chunk:
                current_chunk += "\n\n" + para
            else:
                current_chunk = para
        else:
            if current_chunk:
                chunks.append(current_chunk)
            # If paragraph itself is too long, split it
            if len(para) > max_chunk_size:
                words = para.split()
                current_chunk = ""
                for word in words:
                    if len(current_chunk) + len(word) + 1 <= max_chunk_size:
                        current_chunk += (" " if current_chunk else "") + word
                    else:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = word
            else:
                current_chunk = para
    
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks


# ===== Document Management =====

async def process_uploaded_file(
    file_content: bytes,
    filename: str,
    content_type: Optional[str] = None
) -> Dict[str, Any]:
    """
    Process an uploaded file and store it.
    
    Args:
        file_content: Raw file bytes
        filename: Original filename
        content_type: MIME type
    
    Returns:
        Document metadata dictionary
    """
    ensure_directories()
    
    # Get file extension
    ext = os.path.splitext(filename)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}. Supported types: {', '.join(SUPPORTED_EXTENSIONS.keys())}")
    
    # Check file size
    if len(file_content) > MAX_FILE_SIZE:
        raise ValueError(f"File too large. Maximum size is {MAX_FILE_SIZE / (1024*1024):.0f}MB")
    
    # Generate unique ID
    doc_id = str(uuid.uuid4())
    
    # Save file
    file_path = os.path.join(DOCUMENTS_DIR, f"{doc_id}{ext}")
    with open(file_path, 'wb') as f:
        f.write(file_content)
    
    # Extract text or process image
    extracted_text, extraction_metadata = extract_text_from_file(file_path, ext)
    
    # Check if this is an image for vision
    is_vision_image = extraction_metadata.get("is_vision_image", False)
    
    # Create chunks (only for text documents)
    chunks = chunk_text(extracted_text) if not is_vision_image else [extracted_text]
    
    # Create document metadata
    document = {
        "id": doc_id,
        "filename": filename,
        "extension": ext,
        "content_type": content_type or SUPPORTED_EXTENSIONS.get(ext, 'application/octet-stream'),
        "size": len(file_content),
        "file_path": file_path,
        "uploaded_at": datetime.utcnow().isoformat(),
        "extracted_text": extracted_text,
        "chunks": chunks,
        "chunk_count": len(chunks),
        "text_length": len(extracted_text),
        "is_active": True,
        "is_vision_image": is_vision_image,
        "extraction_metadata": extraction_metadata
    }
    
    # Save to registry (don't save base64 in registry to keep it small)
    registry_doc = document.copy()
    if "base64_data" in registry_doc.get("extraction_metadata", {}):
        registry_doc["extraction_metadata"] = {k: v for k, v in registry_doc["extraction_metadata"].items() if k != "base64_data"}
    
    registry = load_document_registry()
    registry["documents"][doc_id] = registry_doc
    save_document_registry(registry)
    
    # Return metadata (without full text for response)
    return {
        "id": doc_id,
        "filename": filename,
        "extension": ext,
        "size": len(file_content),
        "uploaded_at": document["uploaded_at"],
        "chunk_count": len(chunks),
        "text_length": len(extracted_text),
        "is_active": True,
        "is_vision_image": is_vision_image,
        "preview": extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text
    }


def get_document(doc_id: str) -> Optional[Dict[str, Any]]:
    """Get a document by ID."""
    registry = load_document_registry()
    return registry["documents"].get(doc_id)


def get_document_text(doc_id: str) -> Optional[str]:
    """Get the extracted text from a document."""
    doc = get_document(doc_id)
    if doc:
        return doc.get("extracted_text", "")
    return None


def get_document_image_base64(doc_id: str) -> Optional[str]:
    """Get the base64-encoded image data for a vision document."""
    doc = get_document(doc_id)
    if doc and doc.get("is_vision_image"):
        # Re-read and encode the image file
        file_path = doc.get("file_path")
        if file_path and os.path.exists(file_path):
            try:
                img = Image.open(file_path)
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                
                # Resize if needed
                max_size = 2048
                if max(img.size) > max_size:
                    ratio = max_size / max(img.size)
                    new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                    img = img.resize(new_size, Image.LANCZOS)
                
                buffer = io.BytesIO()
                img.save(buffer, format='PNG')
                return base64.b64encode(buffer.getvalue()).decode('utf-8')
            except Exception as e:
                logger.error(f"Error reading image for vision: {e}")
    return None


def list_documents() -> List[Dict[str, Any]]:
    """List all documents (metadata only)."""
    registry = load_document_registry()
    documents = []
    for doc_id, doc in registry["documents"].items():
        documents.append({
            "id": doc["id"],
            "filename": doc["filename"],
            "extension": doc["extension"],
            "size": doc["size"],
            "uploaded_at": doc["uploaded_at"],
            "chunk_count": doc.get("chunk_count", 1),
            "text_length": doc.get("text_length", 0),
            "is_active": doc.get("is_active", True),
            "is_vision_image": doc.get("is_vision_image", False),
            "preview": doc.get("extracted_text", "")[:200] + "..." if len(doc.get("extracted_text", "")) > 200 else doc.get("extracted_text", "")
        })
    # Sort by upload time, newest first
    documents.sort(key=lambda x: x["uploaded_at"], reverse=True)
    return documents


def delete_document(doc_id: str) -> bool:
    """Delete a document and its file."""
    registry = load_document_registry()
    if doc_id not in registry["documents"]:
        return False
    
    doc = registry["documents"][doc_id]
    
    # Delete file
    file_path = doc.get("file_path")
    if file_path and os.path.exists(file_path):
        os.remove(file_path)
    
    # Remove from registry
    del registry["documents"][doc_id]
    save_document_registry(registry)
    
    return True


def toggle_document_active(doc_id: str, is_active: bool) -> bool:
    """Toggle whether a document is active in the conversation context."""
    registry = load_document_registry()
    if doc_id not in registry["documents"]:
        return False
    
    registry["documents"][doc_id]["is_active"] = is_active
    save_document_registry(registry)
    return True


def get_active_documents_context() -> str:
    """Get combined context from all active text documents."""
    registry = load_document_registry()
    context_parts = []
    
    for doc_id, doc in registry["documents"].items():
        if doc.get("is_active", True) and not doc.get("is_vision_image", False):
            filename = doc.get("filename", "Unknown")
            text = doc.get("extracted_text", "")
            if text:
                context_parts.append(f"=== Document: {filename} ===\n{text}")
    
    if context_parts:
        return "\n\n".join(context_parts)
    return ""


def get_active_vision_images() -> List[Dict[str, Any]]:
    """Get all active vision images with their base64 data."""
    registry = load_document_registry()
    images = []
    
    for doc_id, doc in registry["documents"].items():
        if doc.get("is_active", True) and doc.get("is_vision_image", False):
            base64_data = get_document_image_base64(doc_id)
            if base64_data:
                images.append({
                    "id": doc_id,
                    "filename": doc.get("filename", "Unknown"),
                    "base64_data": base64_data,
                    "width": doc.get("extraction_metadata", {}).get("width", 0),
                    "height": doc.get("extraction_metadata", {}).get("height", 0)
                })
    
    return images


def get_active_document_ids() -> List[str]:
    """Get IDs of all active documents."""
    registry = load_document_registry()
    return [doc_id for doc_id, doc in registry["documents"].items() if doc.get("is_active", True)]
